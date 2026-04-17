
from pptx import Presentation

prs = Presentation("投资心态管理.pptx")

print(f"投资心态管理 PPT - 内容预览\n")
print(f"幻灯片总数：{len(prs.slides)}\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"{'='*50}")
    print(f"幻灯片 {i}")
    print(f"{'='*50}")
    
    # 获取所有文本
    all_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            all_text.append(shape.text.strip())
    
    # 显示前 5 个文本元素
    for j, text in enumerate(all_text[:5], 1):
        # 截断长文本
        if len(text) > 80:
            text = text[:77] + "..."
        print(f"  {j}. {text}")
    
    if len(all_text) > 5:
        print(f"  ... 还有 {len(all_text) - 5} 个文本元素")
    print()
