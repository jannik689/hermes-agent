
from pptx import Presentation

prs = Presentation("投资心态管理.pptx")

print(f"幻灯片总数：{len(prs.slides)}\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"=== 幻灯片 {i} ===")
    # 获取标题
    title = slide.shapes.title.text if slide.shapes.title else "(无标题)"
    print(f"标题：{title}")
    
    # 获取文本内容
    text_items = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_items.append(shape.text.strip())
    
    if text_items:
        print(f"文本元素数量：{len(text_items)}")
    print()
