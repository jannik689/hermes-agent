#!/usr/bin/env python3
"""
TRAE 使用分享 PPT 生成器
使用 python-pptx 创建专业的幻灯片
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# 创建演示文稿 (16:9 比例)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色方案 - Sketch-notes 风格 (温暖、手绘感)
COLORS = {
    'cream': RGBColor(0xFF, 0xF8, 0xF0),      # 背景色
    'charcoal': RGBColor(0x2D, 0x37, 0x48),   # 主文本
    'coral': RGBColor(0xFF, 0x6B, 0x6B),      # 强调色 1
    'teal': RGBColor(0x4E, 0xCD, 0xC4),       # 强调色 2
    'golden': RGBColor(0xFF, 0xD9, 0x3D),     # 强调色 3
    'white': RGBColor(0xFF, 0xFF, 0xFF),
}

def set_background(slide, color=COLORS['cream']):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_background(slide)
    
    # 添加主标题
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['charcoal']
    title_p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(4.0), Inches(10.3), Inches(1.0)
    )
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = COLORS['coral']
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # 添加装饰元素 - 顶部横条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.3)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['teal']
    top_bar.line.fill.background()
    
    return slide

def add_content_slide(prs, title, content_items, visual_note=None):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['charcoal']
    
    # 添加内容
    if content_items:
        content_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.0)
        )
        content_tf = content_box.text_frame
        content_tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = content_tf.paragraphs[0]
            else:
                p = content_tf.add_paragraph()
            
            p.text = f"• {item}" if isinstance(item, str) else item
            p.font.size = Pt(22)
            p.font.color.rgb = COLORS['charcoal']
            p.space_after = Pt(14)
    
    # 添加视觉提示框
    if visual_note:
        note_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(10.5), Inches(5.5), Inches(2.5), Inches(1.5)
        )
        note_box.fill.solid()
        note_box.fill.fore_color.rgb = COLORS['golden']
        note_box.line.color.rgb = COLORS['charcoal']
        note_box.line.width = Pt(2)
        
        note_tf = note_box.text_frame
        note_p = note_tf.paragraphs[0]
        note_p.text = visual_note
        note_p.font.size = Pt(14)
        note_p.font.color.rgb = COLORS['charcoal']
        note_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加双列内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['charcoal']
    
    # 左列标题
    left_title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.5), Inches(5.5), Inches(0.8)
    )
    left_title_tf = left_title_box.text_frame
    left_title_p = left_title_tf.paragraphs[0]
    left_title_p.text = left_title
    left_title_p.font.size = Pt(24)
    left_title_p.font.bold = True
    left_title_p.font.color.rgb = COLORS['coral']
    
    # 左列内容
    left_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.3), Inches(5.5), Inches(4.5)
    )
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i == 0:
            p = left_tf.paragraphs[0]
        else:
            p = left_tf.add_paragraph()
        p.text = f"✅ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['charcoal']
        p.space_after = Pt(10)
    
    # 右列标题
    right_title_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.8)
    )
    right_title_tf = right_title_box.text_frame
    right_title_p = right_title_tf.paragraphs[0]
    right_title_p.text = right_title
    right_title_p.font.size = Pt(24)
    right_title_p.font.bold = True
    right_title_p.font.color.rgb = COLORS['teal']
    
    # 右列内容
    right_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(2.3), Inches(5.5), Inches(4.5)
    )
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i == 0:
            p = right_tf.paragraphs[0]
        else:
            p = right_tf.add_paragraph()
        p.text = f"❌ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['charcoal']
        p.space_after = Pt(10)
    
    # 添加分隔线
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.6), Inches(1.5), Inches(0.1), Inches(5.3)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLORS['charcoal']
    divider.line.fill.background()
    
    return slide

def add_comparison_table_slide(prs, title, headers, data):
    """添加对比表格幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['charcoal']
    
    # 创建表格
    rows = len(data) + 1
    cols = len(headers)
    table_x = Inches(1.0)
    table_y = Inches(1.5)
    table_width = Inches(11.3)
    table_height = Inches(5.0)
    
    table = slide.shapes.add_table(
        rows, cols, table_x, table_y, table_width, table_height
    ).table
    
    # 设置列宽
    column_widths = [Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)]
    for i, width in enumerate(column_widths[:cols]):
        table.columns[i].width = width
    
    # 填充表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['teal']
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLORS['white']
            paragraph.alignment = PP_ALIGN.CENTER
    
    # 填充数据
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = COLORS['charcoal']
                paragraph.alignment = PP_ALIGN.CENTER
    
    return slide

def add_back_cover_slide(prs, title, subtitle, steps):
    """添加封底幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # 添加主标题
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.5)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['coral']
    title_p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.0)
    )
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(26)
    subtitle_p.font.color.rgb = COLORS['charcoal']
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # 添加步骤
    steps_box = slide.shapes.add_textbox(
        Inches(3.0), Inches(4.0), Inches(7.3), Inches(2.5)
    )
    steps_tf = steps_box.text_frame
    steps_tf.word_wrap = True
    
    for i, step in enumerate(steps, 1):
        if i == 1:
            p = steps_tf.paragraphs[0]
        else:
            p = steps_tf.add_paragraph()
        p.text = f"{i}. {step}"
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['charcoal']
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)
    
    # 底部装饰条
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.2), Inches(13.333), Inches(0.3)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = COLORS['golden']
    bottom_bar.line.fill.background()
    
    return slide

# ============ 创建幻灯片 ============

print("创建 TRAE 使用分享 PPT...")

# Slide 1: 封面
print("  Slide 1: 封面")
add_title_slide(prs, "TRAE 使用分享", "AI 编程助手实践指南")

# Slide 2: 什么是 TRAE
print("  Slide 2: 什么是 TRAE")
add_content_slide(prs, "什么是 TRAE？", [
    "智能 AI 编程助手",
    "代码生成与补全",
    "代码理解与解释",
    "智能调试与问题诊断",
    "文档自动生成",
])

# Slide 3: 为什么选择 TRAE
print("  Slide 3: 为什么选择 TRAE")
add_content_slide(prs, "为什么选择 TRAE？", [
    "🚀 提升开发效率 60%",
    "📚 降低学习成本",
    "💎 提高代码质量",
])

# Slide 4: 快速开始
print("  Slide 4: 快速开始")
add_content_slide(prs, "快速开始", [
    "1. 访问官网下载 TRAE",
    "2. 安装对应 IDE 插件",
    "3. 配置 API Key",
    "4. 开始使用",
])

# Slide 5: 代码生成
print("  Slide 5: 代码生成")
add_content_slide(prs, "核心功能：代码生成", [
    "描述需求 → 生成代码",
    "创建 REST API 端点",
    "实现数据验证逻辑",
    "生成单元测试用例",
])

# Slide 6: 代码解释
print("  Slide 6: 代码解释")
add_content_slide(prs, "核心功能：代码解释", [
    "选中代码 → 理解逻辑",
    "阅读他人代码",
    "理解复杂算法",
    "学习新框架源码",
])

# Slide 7: 智能调试
print("  Slide 7: 智能调试")
add_content_slide(prs, "核心功能：智能调试", [
    "错误信息 → 修复建议",
    "错误根因分析",
    "修复方案提供",
    "预防措施建议",
])

# Slide 8: 代码重构
print("  Slide 8: 代码重构")
add_content_slide(prs, "核心功能：代码重构", [
    "选择代码 → 优化结果",
    "性能优化",
    "可读性提升",
    "架构改进",
    "设计模式应用",
])

# Slide 9: Web 开发案例
print("  Slide 9: Web 开发案例")
add_content_slide(prs, "实战案例：Web 开发", [
    "构建用户认证系统",
    "生成 JWT 认证模板",
    "实现密码加密逻辑",
    "创建中间件验证",
    "时间节省：约 60%",
])

# Slide 10: 最佳实践
print("  Slide 10: 最佳实践")
add_two_column_slide(
    prs,
    "最佳实践",
    "✅ 推荐做法",
    ["清晰描述需求", "提供上下文信息", "审查生成代码"],
    "❌ 避免做法",
    ["模糊的需求描述", "完全依赖 AI 代码", "忽略代码审查"],
)

# Slide 11: 提示词技巧
print("  Slide 11: 提示词技巧")
add_content_slide(prs, "提示词技巧", [
    "高效提示公式：[角色] + [任务] + [约束] + [示例]",
    "",
    "示例：",
    "\"作为资深 Python 开发者，创建异步 HTTP 客户端，",
    "使用 aiohttp 库，包含重试机制和超时处理...\"",
])

# Slide 12: 安全与隐私
print("  Slide 12: 安全与隐私")
add_content_slide(prs, "安全与隐私", [
    "🔒 不上传敏感代码",
    "🔒 不输入 API 密钥",
    "🔒 不分享业务逻辑细节",
    "🔒 审查生成代码安全性",
])

# Slide 13: 工具对比
print("  Slide 13: 工具对比")
headers = ["特性", "TRAE", "Copilot", "Cursor"]
data = [
    ["代码生成", "✅", "✅", "✅"],
    ["代码解释", "✅", "⚠️", "✅"],
    ["智能调试", "✅", "❌", "⚠️"],
    ["中文支持", "✅", "⚠️", "✅"],
]
add_comparison_table_slide(prs, "与其他工具对比", headers, data)

# Slide 14: 总结
print("  Slide 14: 总结")
add_content_slide(prs, "总结", [
    "🚀 提升开发效率",
    "📚 降低学习门槛",
    "💡 激发创新思路",
    "🎯 提高代码质量",
    "",
    "关键：AI 是助手，不是替代",
])

# Slide 15: 封底
print("  Slide 15: 封底")
add_back_cover_slide(
    prs,
    "开始你的 TRAE 之旅！",
    "让 AI 成为你的编程伙伴",
    ["下载安装", "配置环境", "尝试第一个任务", "分享你的经验"],
)

# 保存 PPT
output_path = "slide-deck/trae-intro/TRAE 使用分享.pptx"
prs.save(output_path)
print(f"\n✅ PPT 创建完成：{output_path}")
print(f"   共 15 张幻灯片")
